"""
更新2026年7月第三周周报PPT
从两份源文件提取最新数据，汇入周报模板
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

SRC = "/home/hermes/workspace/source/数智工程交付部20260719周报.pptx"
OUT = SRC  # modify in place

prs = Presentation(SRC)

def get_cell(table, row_idx, col_idx):
    """Get cell, handling merged cells gracefully"""
    try:
        return table.cell(row_idx, col_idx)
    except:
        return None

def set_cell_text(cell, text):
    """Set cell text, preserving formatting"""
    if cell is None:
        return
    tf = cell.text_frame
    # Clear existing paragraphs
    for p in tf.paragraphs:
        for r in p.runs:
            r.text = ''
    # Set first paragraph
    tf.paragraphs[0].text = text

def get_table_on_slide(slide):
    """Find the first table on a slide"""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None

def find_textbox(slide, text_contains):
    """Find a text box/auto shape containing text"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if text_contains in para.text:
                    return shape
    return None

def set_textbox_content(shape, text_lines):
    """Replace content of a text box/auto shape with new lines"""
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    # Set first para
    tf.paragraphs[0].text = text_lines[0] if text_lines else ''
    # Add remaining
    for line in text_lines[1:]:
        p = tf.add_paragraph()
        p.text = line

# ============================================================
# 1. Update Slide 4 - 短板总清单
# ============================================================
slide4 = prs.slides[3]  # 0-indexed
table = get_table_on_slide(slide4)
if table:
    print(f"Slide 4 table: {len(table.rows)} rows x {len(table.columns)} cols")
    
    # --- 黄定斌 data mapping (template row index -> new content) ---
    # Each tuple: (row_idx, [col1_new, col2_new, ...]) or None to skip
    # Columns: 0=序号, 1=问题清单, 2=问题来源, 3=完成人, 4=方案措施, 5=完成时间, 6=验收人, 7=状态
    
    updates = {
        # Row 6 = 短板6 (template) = 基站租赁未租 (黄定斌 file 短板7)
        6: {
            4: "未租赁8个站点：1个光伏待整改，1个光伏整改完成，7月22日前完成市电审定单提租赁；1个换接火点造价变更跟移动沟通中；4个网运站点第二步流程曾义成审批中；1个网运站点7月23日前县公司提交资料给曾义成提上会。",
            5: "7月23日",
        },
        # Row 7 = 短板7 = 基站租赁到期开票 (黄定斌 file 短板8)
        7: {
            4: "143站点考核表制作已完成（65座共享，78座独享）。曾义成审核76座独享站(2座考核金额较大未提交审核，已告知国动，待确认处理方案)。36座独享站考核表纸质版7月14日给曾义成签字盖章，40座待曾义成审核。",
            5: "持续推进",
        },
        # Row 8 = 短板8 = 基站巡检 (黄定斌 file 短板9)
        8: {
            4: "累计完成巡检7个完成率2%（7/292）。7月任务61个，完成7个（11.5%）。各县完成率：扶绥0%、宁明2%、大新0%、龙州6%、凭祥0%、江州3%、天等5%。措施：每周通报中心巡检进场，督促加快巡检，7月30日前完成上半年收尾。",
            5: "7月30日",
        },
        # Row 9 = 短板9 = 基站动环监控 (黄定斌 file 短板10)
        9: {
            4: "已收到物联网卡84张，动环材料总价9.52万元，需签采购合同付30%预付款2.86万元后厂家生产，约25天交付。待确认签合同及支付30%预付款。",
            5: "7月30日",
        },
        # Row 10 = 短板10 = 基站电费 (黄定斌 file 短板12)
        10: {
            4: "累计垫付153.5万元，仅回款2.1万元（挂预收无法核销）。王月剑全面负责（黄定斌配合）电费全流程闭环整改。已梳理满足付款条件18个电表44.33万元（2个独享、16个共享）。需中心钳流61个站点已完成60个完成率98%。",
            5: "持续推进",
        },
        # Row 16 = 短板16 = 铁体一体化后评估 (黄定斌 file 短板16 - 已完成)
        16: {
            4: "3月后评估81.6分，沟通铁塔加分，4月加5分至86.6分，5月加分5分至91.6分。已完成。",
            5: "6月30日",
            7: "完成",
        },
        # Row 17 = 短板17 = 低成本乡村振兴 (黄定斌 file 短板15)
        17: {
            4: "涉及3个县份共26个站点（天等17个、龙州3个、大新6个），已跟各县份沟通帮勘察。",
            5: "6月30日",
        },
        # Row 18 = 短板18 = 基站蓄电池低续航 (黄定斌 file 短板13)
        18: {
            4: "112个蓄电池低续航，完成测试62个（55%）。凭祥、龙州、扶绥已完成100%；江州21%、大新9%、天等33%后三。经测试62个站原因：无开关电源2个，无电源模块3个，开关电源坏1个，蓄电池坏56个。已申请通过27套蓄电池。",
            5: "持续推进",
        },
        # Row 19 = 短板19 = 基站备品备件 (黄定斌 file 短板14)
        19: {
            4: "6月10日第一批备品备件到货，6月22日第二批到货，7月17日第三批备件到货，备件仍存在缺口，继续跟国动申请备件补齐缺口。",
            5: "持续推进",
        },
    }
    
    # Also update 黄福华's existing row (Row 12 = 短板12)
    # 黄福华 file has more detailed info on 开完工
    updates[12] = {
        4: "DICT系统开工应完23个（区公司通报22个）：作废13个，业主原因未签6个，还在走合同流程2个，旧项目无法补录1个，属于工程项目已在DICT非正常结1个。智慧工程本月需点开工2个已点1个，需点完工6个（点完工影响财务报账暂未点）。智慧工程需上传资料7个已完成1个，需整改12个已完成2个。",
        5: "7月31日",
    }
    
    # Apply updates
    for row_idx, col_updates in updates.items():
        for col_idx, new_text in col_updates.items():
            cell = get_cell(table, row_idx, col_idx)
            if cell:
                old = cell.text.strip()[:50]
                set_cell_text(cell, new_text)
                print(f"  Row {row_idx}, Col {col_idx}: '{old}...' -> updated ✓")
    
    # Also update the summary line above the table (auto_shape "内容占位符 4")
    for shape in slide4.shapes:
        if shape.has_text_frame and "截至" in shape.text_frame.paragraphs[0].text:
            shape.text_frame.paragraphs[0].text = "截至7月19日短板17个，完成6个；其中一季经分会2个完成1个，条线周例会13个完成5个，移动通报2个完成0个。"
            print("  Summary line updated ✓")
            break
else:
    print("ERROR: No table found on slide 4")

# ============================================================
# 2. Update Slide 8 - 基站业务分析 (黄定斌)
# ============================================================
slide8 = prs.slides[7]  # 0-indexed
for shape in slide8.shapes:
    if shape.has_text_frame and "问题1" in shape.text_frame.paragraphs[0].text:
        new_text = (
            "短板6：基站租赁未租8站（占全区16.33%），租金应开未开票143站1104万。\n"
            "分析及措施：\n"
            "问题1：基站租赁未租赁8个站点：1个光伏待整改，1个光伏整改完成7月22日前完成市电审定单提租赁；"
            "1个换接火点造价变更函移动审核中；4个网运站点第二步流程曾义成审批中；"
            "1个网运站点7月23日前县公司提交资料给曾义成提上会。\n"
            "问题2：143站点考核表制作已完成（65座共享，78座独享）。"
            "曾义成审核76座独享站（2座考核金额较大未提交审核，已告知国动待确认处理方案）。"
            "36座独享站考核表纸质版7月14日给曾义成签字盖章，40座待曾义成审核。\n"
            "问题3：动环84个未安装，动环应用率不高，均靠代维提醒，在线率不可控。"
            "已收到物联网卡84张，动环材料总价9.52万元，需签采购合同付30%预付款2.86万元后厂家生产约25天交付。"
            "待确认签合同及支付30%预付款。\n"
            "进度：光伏1个待整改，1个整改完成；换接火点造价变更与移动沟通中；"
            "4个网运站点曾义成审批中；1个网运站点7月23日前提交资料上会。"
            "独享站考核表36座已签字盖章，40座待审核。"
        )
        # Clear and set
        tf = shape.text_frame
        tf.paragraphs[0].text = new_text
        print("Slide 8 (基站业务) updated ✓")
        break

# ============================================================
# 3. Update Slide 9 - 基站业务分析2 (黄定斌 - 巡检/电费/蓄电池)
# ============================================================
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if shape.has_text_frame and "巡检" in shape.text_frame.paragraphs[0].text:
        new_text = (
            "短板：基站巡检率2%（285站未开始），电费垫付153.5万回款严重滞后，蓄电池低续航112站。\n"
            "分析及措施：\n"
            "问题1：基站巡检累计完成率2%（7/292），7月任务61个仅完成7个。"
            "扶绥、凭祥、宁明前三（相对），龙州、天等、大新后三。"
            "上半年巡检任务未完成，国动将停止备件供应。\n"
            "问题2：电费累计垫付153.5万元，仅回款2.1万元（挂预收无法核销）。"
            "王月剑全面负责（黄定斌配合）电费对账、报账流程跟进、协议催签全流程。"
            "已梳理满足付款条件18个电表44.33万元，需钳流61个站点已完60个（98%）。\n"
            "问题3：蓄电池低续航112个站，测试完成62个（55%）。"
            "凭祥、龙州、扶绥已完成100%；江州21%、大新9%、天等33%严重滞后。"
            "62个站已发邮件申请配件，已通过27套蓄电池。\n"
            "措施：每周通报巡检进场，7月30日前完成上半年收尾。"
            "优先推进直供电快速结算，9个电表27.93万元。"
            "督促各县份加快蓄电池放电测试，特别是江州、大新、天等。"
        )
        tf = shape.text_frame
        tf.paragraphs[0].text = new_text
        print("Slide 9 (基站业务2) updated ✓")
        break

# ============================================================
# 4. Update Slide 13 - 黄福华（低成本及其他）
# ============================================================
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape.has_text_frame and "黄福华" in shape.text_frame.paragraphs[0].text:
        new_text = (
            "短板8（黄福华）：施工项目开完工及时率低（开工25%后三，完工0%末位）；"
            "DICT系统开工项目排名靠后；四川中移传输管线完工排名靠前。\n"
            "分析及措施：\n"
            "问题1：DICT系统开工应完23个（区公司通报22个），其中作废13个，业主原因未签6个，"
            "合同流程中2个，旧项目无法补录1个，属工程项目已非正常结1个。系统及业主原因暂时无法推进，需协调区公司不考核。\n"
            "问题2：智慧工程系统开完工落后。本月需点开工2个已点1个，需点完工6个（7月31日前），"
            "但目前点完工会影响财务报账，暂未点完工。\n"
            "问题3：智慧工程系统工程资料上传。本月需完成上传工程项目7个已完成1个，"
            "需整改工程项目12个已完成2个，各项目分工明确需在7月31日前完成。\n"
            "问题4：四川中移传输管线2025年站点接入一阶段光缆工程项目完工排名靠前。"
            "区公司周例会时未完成16个，截至目前已完成7个，未完成9个。"
            "跟进施工队派工情况，督促施工项目经理提交完工申请。"
        )
        tf = shape.text_frame
        tf.paragraphs[0].text = new_text
        print("Slide 13 (黄福华) updated ✓")
        break

# ============================================================
# 5. Update Slide 15 - 中心通报 (蓄电池 + 巡检 + 纤芯修复)
# ============================================================
slide15 = prs.slides[14]

# Update 蓄电池 table (first table on slide 15)
tables_found = 0
for shape in slide15.shapes:
    if shape.has_table:
        table = shape.table
        if tables_found == 0:
            # 蓄电池 table - update 合计 row
            # Row 8 = 合计 row (after header + 7 districts)
            if len(table.rows) >= 9:
                # Update 合计 row: 需测试数=112, 测试完成数=62, 完成率=55%
                set_cell_text(get_cell(table, 8, 1), "112")   # 需测试数
                set_cell_text(get_cell(table, 8, 2), "62")    # 测试完成数
                set_cell_text(get_cell(table, 8, 3), "55%")   # 完成率
                set_cell_text(get_cell(table, 8, 5), "62")    # 流程申请数
                set_cell_text(get_cell(table, 8, 6), "27")    # 申请成功数
                
                # Update per-district data
                district_updates = {
                    1: ("江州", 24, 5, "21%", 5, 3),
                    2: ("扶绥", 29, 29, "100%", 29, 12),
                    3: ("龙州", 9, 9, "100%", 9, 7),
                    4: ("宁明", 20, 14, "70%", 14, 4),
                    5: ("凭祥", 1, 1, "100%", 1, 0),
                    6: ("大新", 23, 2, "9%", 2, 1),
                    7: ("天等", 6, 2, "33%", 2, 0),
                }
                for r, (name, need, done, rate, app, succ) in district_updates.items():
                    set_cell_text(get_cell(table, r, 1), str(need))
                    set_cell_text(get_cell(table, r, 2), str(done))
                    set_cell_text(get_cell(table, r, 3), rate)
                    set_cell_text(get_cell(table, r, 4), str(r))  # ranking
                    set_cell_text(get_cell(table, r, 5), str(app))
                    set_cell_text(get_cell(table, r, 6), str(succ))
                
                print("  Battery table updated ✓")
            tables_found += 1
        elif tables_found == 1:
            # 巡检 table (second table on slide 15)
            if len(table.rows) >= 5:
                # Update 合计 row
                set_cell_text(get_cell(table, 4, 8), "292")   # 总数
                set_cell_text(get_cell(table, 4, 1), "0")     # 扶绥完成
                set_cell_text(get_cell(table, 4, 2), "1")     # 宁明完成
                set_cell_text(get_cell(table, 4, 3), "0")     # 大新完成
                set_cell_text(get_cell(table, 4, 4), "3")     # 龙州完成
                set_cell_text(get_cell(table, 4, 5), "0")     # 凭祥完成
                set_cell_text(get_cell(table, 4, 6), "1")     # 江州完成
                set_cell_text(get_cell(table, 4, 7), "2")     # 天等完成
                set_cell_text(get_cell(table, 4, 8), "7")     # 合计完成
                set_cell_text(get_cell(table, 4, 1), "0%")    # Update first cell of 完成率 row
                
                # Actually 巡检 table has structure: Row 0=header, Row 1=总数, Row 2=完成数, Row 3=完成率, Row 4=排名
                # Let me re-read: it's 5 rows x 9 cols
                # Row 1 = 总数 row (heads: 区域 | 扶绥 | 宁明 | 大新 | 龙州 | 凭祥 | 江州区 | 天等 | 合计)
                # Row 2 = 完成数 row
                # Row 3 = 完成率 row
                # Row 4 = 排名 row
                
                # Update 总数 (should stay 292)
                # Update 完成数 (Row 2)
                set_cell_text(get_cell(table, 2, 1), "0")    # 扶绥
                set_cell_text(get_cell(table, 2, 2), "1")    # 宁明
                set_cell_text(get_cell(table, 2, 3), "0")    # 大新
                set_cell_text(get_cell(table, 2, 4), "3")    # 龙州
                set_cell_text(get_cell(table, 2, 5), "0")    # 凭祥
                set_cell_text(get_cell(table, 2, 6), "1")    # 江州
                set_cell_text(get_cell(table, 2, 7), "2")    # 天等
                set_cell_text(get_cell(table, 2, 8), "7")    # 合计
                
                # Update 完成率 (Row 3)
                set_cell_text(get_cell(table, 3, 1), "0%")
                set_cell_text(get_cell(table, 3, 2), "2%")
                set_cell_text(get_cell(table, 3, 3), "0%")
                set_cell_text(get_cell(table, 3, 4), "6%")
                set_cell_text(get_cell(table, 3, 5), "0%")
                set_cell_text(get_cell(table, 3, 6), "3%")
                set_cell_text(get_cell(table, 3, 7), "5%")
                set_cell_text(get_cell(table, 3, 8), "2%")
                
                # Update 排名 (Row 4)
                set_cell_text(get_cell(table, 4, 1), "5")
                set_cell_text(get_cell(table, 4, 2), "4")
                set_cell_text(get_cell(table, 4, 3), "5")
                set_cell_text(get_cell(table, 4, 4), "1")
                set_cell_text(get_cell(table, 4, 5), "5")
                set_cell_text(get_cell(table, 4, 6), "3")
                set_cell_text(get_cell(table, 4, 7), "2")
                
                print("  Inspection table updated ✓")
            tables_found += 1
        elif tables_found == 2:
            # 纤芯修复 table (third table)
            # Keep as-is unless we have new data
            tables_found += 1

# Update the text boxes next to tables
for shape in slide15.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.paragraphs[0].text
        if "整体情况：应测" in text:
            shape.text_frame.paragraphs[0].text = "整体情况：应测112个，完成测试62个，测试完成率55%（62/112）。"
            shape.text_frame.paragraphs[1].text = "上周情况：上个周期完成放电测试0个，各县整体进度缓慢。"
            shape.text_frame.paragraphs[2].text = "测试进度：扶绥、龙州、凭祥前三；江州、天等、大新后三。"
            shape.text_frame.paragraphs[3].text = "更换进度：申请通过27套蓄电池，已陆续配送更换中。"
            print("  Battery textbox updated ✓")
        elif "整体情况：累计完成" in text and "巡检" in text:
            shape.text_frame.paragraphs[0].text = "整体情况：7月累计完成巡检7个，完成率2%（7/292）。"
            shape.text_frame.paragraphs[1].text = "上周情况：上周完成巡检2个（龙州2个）。"
            shape.text_frame.paragraphs[2].text = "巡检进度：龙州、天等、江州前三，扶绥、凭祥、大新后三。"
            shape.text_frame.paragraphs[3].text = "巡检要求：上半年巡检任务未完成，国动将对未完成地市停止备件供应，请各中心7月30日前完成巡视。"
            print("  Inspection textbox updated ✓")

# ============================================================
# 6. Update Slide 18 - Appendix 区公司通报
# ============================================================
# Check if 区公司通报 slide needs updating - the existing data already includes 黄福华's items
# Let's verify the data is consistent

# ============================================================
# 7. Add new slide for 黄福华 additional短falls (slides 8-13 area)
# ============================================================
# Add a new slide after slide 13 for 黄福华's additional analysis
# Using the "统一模板" layout which is index 1
slide_layout = prs.slide_layouts[1]  # 统一模板

# Add new slide for 黄福华 additional items
new_slide = prs.slides.add_slide(slide_layout)

# Set title
for shape in new_slide.shapes:
    if shape.has_text_frame and shape.name == "标题 70":
        shape.text_frame.paragraphs[0].text = "二、短板指标分析及措施-施工项目开完工（责任人：黄福华）"
        break

# Find the content placeholder or auto shape
for shape in new_slide.shapes:
    if shape.has_text_frame and shape.name == "内容占位符 4":
        new_text = (
            "短板12：施工项目开完工及时率低（开工25%后三，完工0%末位）。\n"
            "分析及措施：\n"
            "问题1：区公司通报崇左施工项目开工及时率25%（后三：贺州20%、崇左25%、南宁45%），"
            "完工及时率0%（末位）。\n"
            "措施1：从智慧工程系统导出需开完工站点，在规定时间内点开工、完工。\n"
            "进度：本月需点开工2个，已点1个（50%）；需点完工6个（7月31日前），"
            "但点完工会影响财务报账，暂未点完工。\n\n"
            "短板14（黄福华）：政企部门DICT系统开工地市排名靠后。\n"
            "问题1：DICT系统项目排名靠后，需点项目开工提高排名。\n"
            "措施1：由于系统及业主原因，目前暂时无法推进该工作，需协调区公司不考核。\n"
            "进度：DICT系统开工应完23个（区公司通报22个），其中："
            "作废13个，业主原因未签6个，还在走合同流程2个，旧项目无法补录1个，"
            "属于工程项目已在DICT非正常结1个。\n\n"
            "短板16（黄福华）：智慧工程系统工程资料上传情况。\n"
            "问题1：周例会上智慧工程系统上传工程资料。\n"
            "措施1：各项目已分工明确，需7月31前完成该项工作。\n"
            "进度：本月需完成上传工程项目7个，目前已完成1个上传；"
            "本月需整改工程项目12个，已完成2个整改。\n\n"
            "短板17（黄福华）：四川中移传输管线2025年站点接入一阶段光缆工程项目完工排名靠后。\n"
            "问题1：项目完工率排名靠前（区公司要求加快）。\n"
            "措施1：跟进施工队派工情况，督促施工项目经理提交完工申请。\n"
            "进度：区公司周例会时未完成16个，截至目前已完成7个，未完成9个。"
        )
        shape.text_frame.paragraphs[0].text = new_text
        print("New slide (黄福华 additional) added ✓")
        break

# ============================================================
# 8. Move the new slide to correct position (after slide 13, before 后评估)
# ============================================================
# python-pptx doesn't easily reorder slides. 
# The new slide is at the end. We'll note this for the user.

# Save
prs.save(SRC)
print(f"\n✅ Saved: {SRC}")
print(f"Total slides: {len(prs.slides)}")
